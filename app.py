# Streamlit app to convert pdf to ppt
import tempfile
import fitz
from pptx import Presentation
from pptx.util import Mm
from tempfile import TemporaryFile
from PIL import Image
from pathlib import Path
import streamlit as st
import os
import base64


def main():
    st.title("PDF to PPT file converter")
    st.markdown('---')

    st.write("Step 1")
    st.header("Upload your .pdf files here")
    uploaded_files = st.file_uploader("", type=['pdf'], accept_multiple_files=True)

    st.markdown('---')

    st.write("Step 2")
    st.header("Click the button to generate .pptx")

    if st.button("Generate PPTs"):
        if len(uploaded_files)==0:
            st.write("Please upload some files first")

        else:
            download_files_list = []
            for uploaded_file in uploaded_files:
                with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
                    fp = Path(tmp_file.name)
                    fp.write_bytes(uploaded_file.getvalue())
                    img_list = pdf_to_images(tmp_file.name)
                    gen_ppt_from_img_list(img_list, uploaded_file.name[0:-4])
                    download_files_list.append(get_binary_file_downloader_html(uploaded_file.name[0:-4]+'.pptx', uploaded_file.name[0:-4]+'.pptx'))
            st.markdown('---')
            st.write("Step 3")
            st.header("Download your files")
            for link in download_files_list:
                st.markdown(link, unsafe_allow_html=True)

    st.markdown('---')


def pdf_to_images(pdf_filename):
    doc = fitz.open(pdf_filename)
    img_list = []
    for i in range(0,doc.page_count):
        page = doc.loadPage(i)
        pix = page.getPixmap()
        img_list.append(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
    return img_list


def gen_ppt_from_img_list(img_list, name):

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for img in img_list:

        aspect_ratio = img.width/img.height
        if aspect_ratio < 1:
            height = 180
            width = int(height * aspect_ratio)
        else:
            width = 244
            height = int(width / aspect_ratio)

        x_offset = int((254 - width) / 2)
        y_offset = int((190.5 - height) / 2)

        fp = TemporaryFile()
        img.save(fp, "PNG")
        fp.seek(0)

        img_path = fp

        slide = prs.slides.add_slide(blank_slide_layout)
        
        pic = slide.shapes.add_picture(img_path, Mm(x_offset), Mm(y_offset), width=Mm(width), height=Mm(height))

        fp.close()

    prs.save(name+'.pptx')

    return None

# Kudos to GokulNC for this downloader function

def get_binary_file_downloader_html(bin_file, file_label='File'):
    with open(bin_file, 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{os.path.basename(bin_file)}">{file_label}</a>'
    return href


if __name__ == '__main__':
    main()
